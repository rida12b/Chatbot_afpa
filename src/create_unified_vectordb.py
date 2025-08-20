import os
import json
import pickle
import faiss
import numpy as np
import logging
from sentence_transformers import SentenceTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tqdm import tqdm
from azure.storage.blob import ContainerClient
import fitz  # PyMuPDF pour PDF
from pptx import Presentation  # Pour PPTX
import io # Pour manipuler les streams en mémoire
from dotenv import load_dotenv
from ingestion import extract_text

# Chemins des dossiers
logger = logging.getLogger(__name__)

# Charger les variables d'environnement
load_dotenv()

def load_local_documents_and_urls():
    """Charge les documents locaux depuis data/documents et les URLs depuis documents/Url_nom_FINA.json.

    Retourne:
      - documents: dict {doc_id: {"content": str, "metadata": {...}}}
      - url_dict: dict {filename: {"url": str, "category": str}}
    """
    documents = {}
    url_dict = {}

    docs_root = os.path.join(DATA_DIR, "documents")
    if not os.path.exists(docs_root):
        logger.warning("Le répertoire local des documents n'existe pas: %s", docs_root)
        return documents, url_dict

    logger.info("Exploration des documents locaux dans %s...", docs_root)
    for root, dirs, files in os.walk(docs_root):
        rel_path = os.path.relpath(root, docs_root)
        category = "general_local" if rel_path == '.' else rel_path
        for file in files:
            file_path = os.path.join(root, file)
            try:
                text = extract_text(file_path)
                if not text or not isinstance(text, str) or not text.strip():
                    logger.warning("Texte vide ou non extractible pour: %s", file_path)
                    continue
                doc_id = os.path.splitext(file)[0]
                documents[f"local_{doc_id}"] = {
                    "content": text,
                    "metadata": {
                        "title": doc_id,
                        "source": file,
                        "path": file_path,
                        "category": category
                    }
                }
            except Exception:
                logger.exception("Erreur lors de l'extraction locale pour %s", file_path)

    # Charger les URLs locales depuis documents/Url_nom_FINA.json
    urls_json_path = os.path.join(BASE_DIR, "documents", "Url_nom_FINA.json")
    if os.path.exists(urls_json_path):
        try:
            with open(urls_json_path, 'r', encoding='utf-8') as f:
                raw = json.load(f)
            # Attendu: liste d'objets {"Nom": filename, "URL": url, "Categorie": opt}
            for item in raw if isinstance(raw, list) else []:
                if isinstance(item, dict) and "Nom" in item and "URL" in item:
                    url_dict[item["Nom"]] = {"url": item["URL"], "category": item.get("Categorie", "FINA_Local")}
            logger.info("URLs locales chargées: %d", len(url_dict))
        except Exception:
            logger.exception("Erreur lors du chargement de documents/Url_nom_FINA.json")
    else:
        logger.warning("Fichier d'URLs locales introuvable: %s", urls_json_path)

    return documents, url_dict
BASE_DIR = os.path.dirname(os.path.dirname(__file__))
DATA_DIR = os.path.join(BASE_DIR, "data")
PROCESSED_DIR = os.path.join(DATA_DIR, "processed")  # Dossier contenant les documents traités par ingestion.py
MODELS_DIR = os.path.join(BASE_DIR, "models")
URL_METADATA_FILE = os.path.join(DATA_DIR, "vectors", "metadata.json")

# Fichiers de sortie
UNIFIED_INDEX_FILE = os.path.join(MODELS_DIR, "unified_index.bin")
UNIFIED_CHUNKS_FILE = os.path.join(MODELS_DIR, "unified_chunks.json")
UNIFIED_METADATA_FILE = os.path.join(MODELS_DIR, "unified_metadata.pkl")

# Paramètres
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"

# --- Fonctions d'extraction de texte (copiées/adaptées depuis ingestion.py) ---

def extract_text_from_pdf_stream(blob_name, blob_stream):
    """Extrait le texte d'un flux de données PDF."""
    try:
        # Lire tout le contenu du stream Azure en mémoire d'abord
        pdf_data = blob_stream.readall()
        # Ensuite, passer ces bytes à fitz.open
        doc = fitz.open(stream=pdf_data, filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text.strip()
    except Exception as e:
        logger.error("Erreur lors de l'extraction du PDF (stream) %s: %s", blob_name, str(e))
        return None

def extract_text_from_pptx_stream(blob_name, blob_stream):
    """Extrait le texte d'un flux de données PPTX."""
    try:
        # Presentation a besoin d'un objet fichier-like, BytesIO fait l'affaire
        prs = Presentation(io.BytesIO(blob_stream.readall()))
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text)
        return text.strip()
    except Exception as e:
        logger.error("Erreur lors de l'extraction du PPTX (stream) %s: %s", blob_name, str(e))
        return None

def extract_text_from_blob(blob_name, blob_client):
    """
    Télécharge un blob et extrait son contenu textuel en fonction de son extension.
    """
    content_text = None
    file_extension = blob_name.split('.')[-1].lower()
    
    logger.info("Téléchargement du blob: %s (type: %s)", blob_name, file_extension)
    blob_stream = blob_client.download_blob() # Renvoie un StorageStreamDownloader

    if file_extension == 'pdf':
        content_text = extract_text_from_pdf_stream(blob_name, blob_stream)
    elif file_extension == 'pptx':
        content_text = extract_text_from_pptx_stream(blob_name, blob_stream)
    elif file_extension == 'txt': # Garder la logique pour les .txt directs
        try:
            content_text = blob_stream.readall().decode('utf-8')
        except Exception as e:
            logger.error("Erreur lors de la lecture du blob texte %s: %s", blob_name, str(e))
            content_text = None
    else:
        logger.warning("Type de blob non supporté pour extraction directe: %s", blob_name)
        # On pourrait ajouter ici d'autres extracteurs si nécessaire

    if content_text:
        logger.info("Texte extrait de %s", blob_name)
    else:
        logger.warning("Échec de l'extraction de texte pour %s", blob_name)
        
    return content_text

# --- Fin des fonctions d'extraction ---

def load_documents(azure_sas_url=None):
    """Recense les éléments à traiter sans charger leur contenu.

    Retourne une liste d'entrées:
    - Local: {"type": "local", "path": file_path, "name": file, "category": category}
    - Azure: {"type": "azure", "blob_client": blob_client, "name": blob.name, "category": category}
    """
    files_to_process = []
    logger.info("Recensement des documents à traiter...")

    # Azure
    if azure_sas_url:
        logger.info("Connexion à Azure Blob Storage avec l'URL SAS...")
        try:
            container_client = ContainerClient.from_container_url(azure_sas_url)
            azure_found = 0
            logger.info("Exploration du conteneur Azure...")
            for blob in container_client.list_blobs():
                ext = os.path.splitext(blob.name)[1].lower()
                if ext not in [".pdf", ".pptx", ".txt"]:
                    continue
                category_parts = blob.name.split('/')[:-1]
                category = "/".join(category_parts) if category_parts else "general_azure"
                files_to_process.append({
                    "type": "azure",
                    "blob_client": container_client.get_blob_client(blob.name),
                    "name": blob.name,
                    "category": category
                })
                azure_found += 1
            logger.info("Blobs éligibles trouvés: %d", azure_found)
        except Exception as e:
            logger.error("Erreur de connexion ou de listing des blobs Azure: %s", str(e))
            logger.warning("Poursuite avec les fichiers locaux uniquement.")

    # Local
    if os.path.exists(PROCESSED_DIR):
        logger.info("Exploration du dossier local %s...", PROCESSED_DIR)
        local_found = 0
        for root, dirs, files in os.walk(PROCESSED_DIR):
            rel_path = os.path.relpath(root, PROCESSED_DIR)
            txt_files = [f for f in files if f.endswith('.txt')]
            if txt_files:
                logger.info("Dossier %s: %d fichiers .txt", rel_path, len(txt_files))
            for file in txt_files:
                file_path = os.path.join(root, file)
                category = "general_local" if rel_path == '.' else rel_path
                files_to_process.append({
                    "type": "local",
                    "path": file_path,
                    "name": file,
                    "category": category
                })
                local_found += 1
        logger.info("Fichiers locaux éligibles trouvés: %d", local_found)
    else:
        logger.warning("Le dossier %s n'existe pas.", PROCESSED_DIR)

    logger.info("Total éléments à traiter: %d", len(files_to_process))
    return files_to_process

def load_url_metadata(url_data_from_azure=None):
    """Charge les métadonnées des URLs à partir des données JSON fournies (depuis Azure)."""
    logger.info("Traitement des métadonnées des URLs fournies...")
    
    if not url_data_from_azure:
        logger.warning("Aucune donnée URL fournie à traiter.")
        return {}

    if not isinstance(url_data_from_azure, list):
        logger.warning("Les données URL fournies ne sont pas une liste comme attendu.")
        return {}

    url_dict = {}
    loaded_count = 0
    for item in url_data_from_azure:
        if isinstance(item, dict) and "Nom" in item and "URL" in item:
            file_name = item["Nom"] 
            url = item["URL"]
            # La clé du dictionnaire sera le nom de fichier exact, ex: "Afpa - Solde à quantité reçue.pdf"
            if file_name not in url_dict:
                url_dict[file_name] = {"url": url, "category": item.get("Categorie", "FINA_Azure")} # Utiliser Categorie si existe, sinon FINA_Azure
                loaded_count += 1
            else:
                logger.warning("Nom de fichier dupliqué dans les données URL ignoré : %s", file_name)
        else:
            logger.warning("Item URL invalide ignoré: %s", item)
            
    if loaded_count > 0:
        logger.info("%d entrées URL traitées et chargées dans url_dict.", loaded_count)
    else:
        logger.warning("Aucune entrée URL valide n'a été chargée depuis les données fournies.")
        
    return url_dict

def create_chunks(documents):
    """Découpe les documents en chunks"""
    logger.info("Découpage des documents en chunks...")
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", " ", ""]
    )
    
    chunks = []
    metadata = []
    
    for doc_id, doc in tqdm(documents.items(), desc="Création des chunks"):
        doc_chunks = text_splitter.split_text(doc["content"])
        
        for i, chunk in enumerate(doc_chunks):
            if chunk.strip():  # Ignorer les chunks vides
                chunks.append(chunk)
                
                # Créer les métadonnées pour ce chunk
                chunk_metadata = doc["metadata"].copy()
                chunk_metadata["chunk_id"] = i
                chunk_metadata["doc_id"] = doc_id
                metadata.append(chunk_metadata)
    
    logger.info("%d chunks créés", len(chunks))
    return chunks, metadata

def associate_urls_to_metadata(metadata, url_dict):
    """Associe les URLs aux métadonnées des chunks en utilisant le nom de fichier source."""
    logger.info("Association des URLs aux chunks...")
    
    associated_count = 0
    for meta_item in tqdm(metadata, desc="Association des URLs"):
        # meta_item["source"] contient le nom du blob, ex: "FINA/Afpa - Solde à quantité reçue.pdf"
        # ou le nom de fichier local, ex: "Afpa - Solde à quantité reçue.pdf"
        # Nous avons besoin du nom de base du fichier pour correspondre avec les clés de url_dict ("Nom" du JSON)
        
        # Extraire le nom de fichier de base de meta_item["source"]
        source_filename = os.path.basename(meta_item["source"])
        
        if source_filename in url_dict:
            url_data = url_dict[source_filename]
            meta_item["url"] = url_data["url"]
            meta_item["url_category"] = url_data["category"]
            associated_count += 1
        # Pas de recherche partielle pour l'instant pour éviter les fausses correspondances avec cette structure
        # Si la clé "Nom" dans le JSON est toujours exacte, la correspondance directe est plus sûre.
    
    logger.info("%d/%d chunks associés à des URLs via le nom de fichier source.", associated_count, len(metadata))
    return metadata

def create_embeddings(chunks):
    """Crée des embeddings pour les chunks"""
    logger.info("Chargement du modèle d'embeddings %s...", EMBEDDING_MODEL)
    model = SentenceTransformer(EMBEDDING_MODEL)
    
    logger.info("Création des embeddings...")
    embeddings = []
    
    # Traiter les chunks par lots
    batch_size = 32
    for i in tqdm(range(0, len(chunks), batch_size), desc="Création des embeddings"):
        batch = chunks[i:i+batch_size]
        batch_embeddings = model.encode(batch)
        embeddings.extend(batch_embeddings)
    
    embeddings = np.array(embeddings).astype('float32')
    logger.info("%d embeddings créés de dimension %d", len(embeddings), embeddings.shape[1])
    
    return embeddings

def build_faiss_index(embeddings):
    """Construit l'index FAISS"""
    logger.info("Construction de l'index FAISS...")
    
    # Créer l'index
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    
    # Ajouter les vecteurs à l'index
    index.add(embeddings)
    
    logger.info("Index FAISS créé avec %d vecteurs", index.ntotal)
    return index

def save_unified_data(index, chunks, metadata):
    """Sauvegarde les données unifiées"""
    logger.info("Sauvegarde des données unifiées...")
    
    # Créer le dossier models s'il n'existe pas
    os.makedirs(MODELS_DIR, exist_ok=True)
    
    # Sauvegarder l'index FAISS
    faiss.write_index(index, UNIFIED_INDEX_FILE)
    logger.info("Index FAISS sauvegardé dans %s", UNIFIED_INDEX_FILE)
    
    # Sauvegarder les chunks
    with open(UNIFIED_CHUNKS_FILE, 'w', encoding='utf-8') as f:
        json.dump({"chunks": chunks}, f, ensure_ascii=False, indent=2)
    logger.info("Chunks sauvegardés dans %s", UNIFIED_CHUNKS_FILE)
    
    # Sauvegarder les métadonnées
    with open(UNIFIED_METADATA_FILE, 'wb') as f:
        pickle.dump(metadata, f)
    logger.info("Métadonnées sauvegardées dans %s", UNIFIED_METADATA_FILE)

def main():
    """Fonction principale"""
    logger.info("Démarrage de la création de la base de données vectorielle unifiée...")
    
    DATA_SOURCE = os.environ.get("DATA_SOURCE", "local").lower()
    AZURE_SAS_URL = os.environ.get("AZURE_SAS_URL")

    parsed_url_json_from_azure = None # Pour stocker les données JSON parsées

    if DATA_SOURCE == 'azure':
        if AZURE_SAS_URL:
            logger.info("Tentative de lecture du fichier JSON d'URLs depuis Azure (FINA/name_url_files.json)...")
            try:
                container_client = ContainerClient.from_container_url(AZURE_SAS_URL)
                json_blob_client = container_client.get_blob_client("FINA/name_url_files.json")
                
                if json_blob_client.exists():
                    downloader = json_blob_client.download_blob()
                    json_content_bytes = downloader.readall()
                    json_content_str = json_content_bytes.decode('utf-8')
                    logger.debug("Contenu du fichier FINA/name_url_files.json (premiers 500 caractères): %s", json_content_str[:500] + ("... (tronqué)" if len(json_content_str) > 500 else ""))
                    
                    try:
                        parsed_url_json_from_azure = json.loads(json_content_str)
                        logger.info("JSON des URLs Azure parsé avec succès.")
                    except json.JSONDecodeError as je:
                        logger.error("Erreur de décodage JSON pour FINA/name_url_files.json: %s", je)
                else:
                    logger.warning("Le fichier FINA/name_url_files.json n'a pas été trouvé sur Azure.")
            except Exception:
                logger.exception("Erreur lors de la tentative de lecture de FINA/name_url_files.json depuis Azure")
        else:
            logger.error("DATA_SOURCE=azure mais AZURE_SAS_URL n'est pas défini")

        # 1. Recenser les éléments à traiter sur Azure
        files_to_process = load_documents(azure_sas_url=AZURE_SAS_URL)
        if not files_to_process:
            logger.error("Aucun document/bloc à traiter trouvé.")
            return

        # 2. Charger/préparer les métadonnées des URLs depuis Azure
        url_dict = load_url_metadata(url_data_from_azure=parsed_url_json_from_azure)
        if not url_dict:
            logger.warning("Aucune métadonnée URL valide chargée. L'association d'URLs sera limitée.")

    else:
        # Mode local
        logger.info("DATA_SOURCE=local: Chargement des documents et URLs locales")
        documents_local, url_dict = load_local_documents_and_urls()
        if not documents_local:
            logger.error("Aucun document local trouvé à traiter.")
            return
        # Convertir le dict de documents locaux en files_to_process-like entries
        files_to_process = []
        for doc_id, doc in documents_local.items():
            files_to_process.append({
                "type": "local",
                "path": doc["metadata"].get("path", doc["metadata"].get("source", "")),
                "name": doc["metadata"].get("source", doc_id + ".txt"),
                "category": doc["metadata"].get("category", "general_local")
            })

    # 3. Préparer les composants partagés
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", " ", ""]
    )
    logger.info("Chargement du modèle d'embeddings %s...", EMBEDDING_MODEL)
    model = SentenceTransformer(EMBEDDING_MODEL)

    index = None
    all_chunks = []
    all_metadata = []

    processed_files = 0
    error_count = 0

    # 4. Traitement incrémental des fichiers/blobs
    for entry in files_to_process:
        file_label = entry.get("path") or entry.get("name")
        try:
            # Extraire le contenu texte
            if entry["type"] == "local":
                with open(entry["path"], 'r', encoding='utf-8') as f:
                    content = f.read()
            else:  # azure
                content = extract_text_from_blob(entry["name"], entry["blob_client"])

            if not content or not content.strip():
                logger.warning("Contenu vide ignoré: %s", file_label)
                continue

            # Métadonnées document
            if entry["type"] == "local":
                doc_id = os.path.splitext(os.path.basename(entry["name"]))[0]
                source = entry["name"]
                path = entry["path"]
            else:
                doc_id = os.path.splitext(os.path.basename(entry["name"]))[0]
                source = entry["name"]
                container_name = getattr(entry["blob_client"], "container_name", "")
                path = f"azure_blob://{container_name}/{entry['name']}" if container_name else f"azure_blob://{entry['name']}"

            category = entry.get("category", "general")

            # Chunking
            doc_chunks = text_splitter.split_text(content)
            if not doc_chunks:
                logger.warning("Aucun chunk généré pour: %s", file_label)
                continue

            # Créer métadonnées des chunks
            file_metadata = []
            for i, chunk in enumerate(doc_chunks):
                if not chunk.strip():
                    continue
                meta = {
                    "title": doc_id,
                    "source": source,
                    "path": path,
                    "category": category,
                    "chunk_id": i,
                    "doc_id": f"{entry['type']}_{doc_id}"
                }
                file_metadata.append(meta)

            # Embeddings pour les chunks de ce fichier (en lot)
            file_chunks = [c for c in doc_chunks if c.strip()]
            file_embeddings = model.encode(file_chunks)
            file_embeddings = np.array(file_embeddings).astype('float32')

            # Initialiser ou alimenter l'index FAISS
            if index is None:
                dimension = file_embeddings.shape[1]
                index = faiss.IndexFlatL2(dimension)
            index.add(file_embeddings)

            # Accumuler pour sauvegarde finale
            all_chunks.extend(file_chunks)
            all_metadata.extend(file_metadata)
            processed_files += 1

            logger.info("Fichier traité: %s | chunks: %d", file_label, len(file_chunks))

        except Exception:
            error_count += 1
            logger.exception("Erreur lors du traitement du fichier %s", file_label)
            continue

    if not all_chunks or index is None:
        logger.error("Aucun chunk ou index non initialisé. Arrêt.")
        return

    # 5. Associer les URLs aux métadonnées (post-traitement)
    all_metadata = associate_urls_to_metadata(all_metadata, url_dict)

    # 6. Sauvegarder les données unifiées
    save_unified_data(index, all_chunks, all_metadata)

    logger.info("Base de données vectorielle unifiée créée avec succès!")
    logger.info("Statistiques finales:")
    logger.info("  - Fichiers traités: %d", processed_files)
    logger.info("  - Erreurs: %d", error_count)
    logger.info("  - Chunks créés: %d", len(all_chunks))
    logger.info("  - Dimension des embeddings: %d", index.d)
    logger.info("  - Taille de l'index FAISS: %d vecteurs", index.ntotal)

if __name__ == "__main__":
    main() 