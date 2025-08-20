import os
import fitz  # PyMuPDF pour PDF
import pytesseract  # OCR pour images
import docx  # Lecture de Word
import pandas as pd  # Lecture d'Excel et CSV
import cv2  # Manipulation d'images pour OCR
from pptx import Presentation  # Lecture de PPTX
import re  # Pour le nettoyage de texte
from datetime import datetime

# 📂 Chemins des dossiers
BASE_DIR = os.path.dirname(os.path.dirname(__file__))
DATA_DIR = os.path.join(BASE_DIR, "data")
DOC_DIR = os.path.join(DATA_DIR, "documents")
PROCESSED_DIR = os.path.join(DATA_DIR, "processed")

# Création des dossiers nécessaires
for directory in [DATA_DIR, DOC_DIR, PROCESSED_DIR]:
    if not os.path.exists(directory):
        os.makedirs(directory)
        print(f"📁 Dossier créé : {directory}")

print(f"🔍 Dossier de données : {DATA_DIR}")
print(f"📂 Dossier documents source : {DOC_DIR}")
print(f"📂 Dossier documents traités : {PROCESSED_DIR}")

print("🚀 Script ingestion.py (TEST) bien exécuté !")

# 📄 **Extraction texte PDF**
def extract_text_from_pdf(file_path):
    """Extrait le texte d'un fichier PDF"""
    try:
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text.strip()
    except Exception as e:
        print(f"❌ Erreur lors de l'extraction du PDF {file_path}: {str(e)}")
        return None

# 📝 **Extraction texte DOCX/DOC**
def extract_text_from_docx(file_path):
    try:
        # Essayer d'abord avec python-docx (pour DOCX)
        try:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs]).strip()
        except Exception as docx_error:
            # Si ça échoue et que c'est un fichier DOC, essayer une approche alternative
            if file_path.lower().endswith('.doc'):
                # Pour les fichiers DOC, on peut essayer de les traiter comme des fichiers binaires
                # et extraire le texte brut
                with open(file_path, 'rb') as f:
                    content = f.read()
                    # Extraction basique de texte ASCII/Unicode
                    text = ""
                    for i in range(0, len(content), 2):
                        if i+1 < len(content):
                            char = content[i:i+2]
                            if 32 <= char[0] <= 126 and char[1] == 0:  # ASCII printable
                                text += chr(char[0])
                    # Nettoyer le texte
                    text = re.sub(r'[^\x20-\x7E\n]', '', text)
                    text = re.sub(r'\s+', ' ', text).strip()
                    return text if text else "[ERREUR DOC] Contenu non extractible"
            else:
                raise docx_error
    except Exception as e:
        return f"[ERREUR DOCX] {file_path} : {e}"

# 📊 **Extraction texte Excel / CSV**
def extract_text_from_excel(file_path):
    try:
        # Essayer différentes options pour Excel
        try:
            if file_path.endswith((".xls", ".xlsx")):
                # Essayer avec engine='openpyxl' pour xlsx
                df = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
            else:
                df = pd.read_csv(file_path)
            
            # Convertir en texte
            if isinstance(df, dict):
                text = "\n\n".join([f"Feuille: {sheet_name}\n" + sheet_df.to_string(index=False) 
                                   for sheet_name, sheet_df in df.items()])
            else:
                text = df.to_string(index=False)
            
            return text.strip()
        except Exception as excel_error:
            # Approche alternative pour les fichiers Excel problématiques
            try:
                # Essayer avec engine='xlrd' pour xls
                if file_path.endswith(".xls"):
                    df = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
                    text = "\n\n".join([f"Feuille: {sheet_name}\n" + sheet_df.to_string(index=False) 
                                       for sheet_name, sheet_df in df.items()])
                    return text.strip()
                else:
                    raise excel_error
            except:
                # Si tout échoue, retourner un message d'erreur plus descriptif
                return f"[ERREUR EXCEL] Impossible de lire {os.path.basename(file_path)} - Format non supporté ou fichier corrompu"
    except Exception as e:
        return f"[ERREUR EXCEL] {file_path} : {e}"

# 📽️ **Extraction texte des images (OCR)**
def extract_text_from_image(file_path):
    try:
        # Lire l'image
        image = cv2.imread(file_path)
        if image is None:
            return f"[ERREUR OCR] Impossible de lire l'image {file_path}"
        
        # Prétraitement pour améliorer l'OCR
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        # Appliquer un seuillage adaptatif
        thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                      cv2.THRESH_BINARY, 11, 2)
        # Débruitage
        denoised = cv2.fastNlMeansDenoising(thresh, None, 10, 7, 21)
        
        # OCR avec Tesseract
        try:
            text = pytesseract.image_to_string(denoised)
            if not text.strip():
                # Si pas de texte, essayer directement sur l'image en niveaux de gris
                text = pytesseract.image_to_string(gray)
            return text.strip()
        except Exception as ocr_error:
            # Si Tesseract n'est pas installé ou échoue
            return f"[ERREUR OCR] {str(ocr_error)}"
    except Exception as e:
        return f"[ERREUR OCR] {file_path} : {e}"

# 📂 **Extraction texte PPTX/PPT**
def extract_text_from_pptx(file_path):
    try:
        # Pour PPTX
        if file_path.lower().endswith('.pptx'):
            prs = Presentation(file_path)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text)
            return text.strip()
        # Pour PPT (ancienne version)
        elif file_path.lower().endswith('.ppt'):
            # Approche alternative pour PPT: extraire du texte brut
            with open(file_path, 'rb') as f:
                content = f.read()
                # Extraction basique de texte ASCII/Unicode
                text = ""
                for i in range(0, len(content)):
                    if 32 <= content[i] <= 126:  # ASCII printable
                        text += chr(content[i])
                # Nettoyer le texte
                text = re.sub(r'[^\x20-\x7E\n]', '', text)
                text = re.sub(r'\s+', ' ', text).strip()
                return text if text else "[ERREUR PPT] Contenu non extractible"
    except Exception as e:
        return f"[ERREUR PPTX/PPT] {file_path} : {e}"

# 📝 **Extraction texte TXT**
def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().strip()
    except Exception as e:
        return f"[ERREUR TXT] {file_path} : {e}"

# **📌 Sélectionne la bonne extraction selon l'extension**
def extract_text(file_path):
    ext = file_path.split('.')[-1].lower()
    print(f"  🔍 Type de fichier détecté : {ext}")

    extraction_methods = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "doc": extract_text_from_docx,
        "xlsx": extract_text_from_excel,
        "xls": extract_text_from_excel,
        "csv": extract_text_from_excel,
        "png": extract_text_from_image,
        "jpg": extract_text_from_image,
        "jpeg": extract_text_from_image,
        "tiff": extract_text_from_image,
        "pptx": extract_text_from_pptx,
        "ppt": extract_text_from_pptx,
        "txt": extract_text_from_txt,
    }

    if ext in extraction_methods:
        return extraction_methods[ext](file_path)
    else:
        return f"[ERREUR] Format non supporté : {ext}"

# Fonction pour sauvegarder le texte extrait
def save_extracted_text(text, source_path, url_mappings=None):
    """Sauvegarde le texte extrait dans un fichier .txt avec les URLs correspondantes"""
    if not text or not text.strip():
        return False
        
    # Créer un nom de fichier basé sur le chemin source
    base_name = os.path.splitext(os.path.basename(source_path))[0]
    txt_path = os.path.join(PROCESSED_DIR, f"{base_name}.txt")
    
    try:
        # Préparer l'en-tête avec les métadonnées
        header = f"SOURCE: {source_path}\n"
        header += f"TITRE: {os.path.basename(source_path)}\n"
        
        # Recherche d'URL avec le système de correspondance sophistiqué
        if url_mappings:
            # Normaliser le titre pour la recherche
            normalized_title = base_name.lower().strip()
            url_found = False
            
            # 1. Recherche exacte
            for name, data in url_mappings.items():
                if name.lower().strip() == normalized_title:
                    header += f"URL: {data['url']}\n"
                    header += f"CATEGORIE: {data.get('category', 'Non spécifiée')}\n"
                    url_found = True
                    print(f"✓ URL trouvée (correspondance exacte) pour {base_name}")
                    break
            
            # 2. Recherche par inclusion si pas de correspondance exacte
            if not url_found:
                best_match = None
                best_score = 0.7  # Seuil minimum de similarité
                
                for name, data in url_mappings.items():
                    name_lower = name.lower().strip()
                    
                    # Vérifier si le titre est contenu dans le nom ou vice versa
                    if normalized_title in name_lower or name_lower in normalized_title:
                        score = 0.85  # Score élevé pour l'inclusion
                        if score > best_score:
                            best_score = score
                            best_match = data
                
                if best_match:
                    header += f"URL: {best_match['url']}\n"
                    header += f"CATEGORIE: {best_match.get('category', 'Non spécifiée')}\n"
                    print(f"✓ URL trouvée (correspondance partielle) pour {base_name}")
                else:
                    print(f"⚠️ Aucune URL trouvée pour {base_name}")
            
        header += f"DATE_EXTRACTION: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        header += "---\n\n"
        
        # Écrire le fichier avec l'en-tête et le contenu
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(header + text)
        return True
    except Exception as e:
        print(f"❌ Erreur lors de la sauvegarde de {txt_path}: {str(e)}")
        return False

# **📂 Parcours tous les fichiers du dossier data/doc/** et sous-dossiers
def process_documents(url_mappings=None):
    """Traite tous les documents dans le dossier doc de manière récursive"""
    print("\n🚀 Démarrage du traitement des documents...")
    
    # Si aucun mapping d'URL n'est fourni, essayer de les charger
    if url_mappings is None:
        try:
            from app import load_url_mappings
            url_mappings = load_url_mappings()
            print(f"✓ {len(url_mappings)} mappings d'URLs chargés")
        except Exception as e:
            print(f"⚠️ Impossible de charger les mappings d'URLs: {e}")
            url_mappings = {}
    
    # Statistiques
    stats = {
        'total': 0,
        'processed': 0,
        'success': 0,
        'formats': {}
    }
    
    # Parcourir récursivement tous les fichiers
    for root, dirs, files in os.walk(DOC_DIR):
        # Créer la structure de dossiers dans processed
        rel_path = os.path.relpath(root, DOC_DIR)
        if rel_path != '.':
            os.makedirs(os.path.join(PROCESSED_DIR, rel_path), exist_ok=True)
        
        for file_name in files:
            file_path = os.path.join(root, file_name)
            print(f"\n📄 Traitement de : {os.path.relpath(file_path, DOC_DIR)}")
            stats['total'] += 1
            
            # Obtenir l'extension
            ext = os.path.splitext(file_name)[1].lower()
            stats['formats'][ext] = stats['formats'].get(ext, 0) + 1
            
            # Extraire le texte
            text = extract_text(file_path)
            if text and not text.startswith('[ERREUR]'):
                stats['processed'] += 1
                # Créer le même chemin relatif dans processed
                processed_dir = os.path.join(PROCESSED_DIR, rel_path)
                os.makedirs(processed_dir, exist_ok=True)
                processed_path = os.path.join(processed_dir, os.path.splitext(file_name)[0] + '.txt')
                
                try:
                    # Sauvegarder avec les URLs
                    if save_extracted_text(text, os.path.relpath(file_path, DOC_DIR), url_mappings):
                        stats['success'] += 1
                        print(f"✅ Texte extrait et sauvegardé dans : {os.path.relpath(processed_path, PROCESSED_DIR)}")
                except Exception as e:
                    print(f"❌ Échec de la sauvegarde : {str(e)}")
            else:
                print(f"❌ Échec de l'extraction : {text}")
    
    # Calculer le taux de succès
    success_rate = (stats['success'] / stats['total'] * 100) if stats['total'] > 0 else 0
    
    # Afficher les statistiques
    print("\n📊 Statistiques de traitement:")
    print(f"  - Documents trouvés: {stats['total']}")
    print(f"  - Documents traités: {stats['processed']}")
    print(f"  - Documents extraits avec succès: {stats['success']}")
    print(f"  - Taux de succès: {success_rate:.1f}%")
    print("\n📝 Formats traités:")
    for fmt, count in stats['formats'].items():
        print(f"  - {fmt}: {count} fichiers")
    
    return stats

# Lancer le traitement si le script est exécuté directement
if __name__ == "__main__":
    print("\n🚀 Démarrage du traitement des documents de test...")
    print(f"📂 Dossier source : {DOC_DIR}")
    print(f"📂 Dossier destination : {PROCESSED_DIR}")
    
    # Vérifier que le dossier source existe et contient des fichiers
    if not os.path.exists(DOC_DIR):
        print(f"❌ Le dossier source {DOC_DIR} n'existe pas!")
        exit(1)
    
    files = os.listdir(DOC_DIR)
    if not files:
        print(f"❌ Le dossier source {DOC_DIR} est vide!")
        exit(1)
    
    print(f"📄 Nombre de fichiers trouvés : {len(files)}")
    print("📝 Liste des fichiers :")
    for file in files:
        print(f"  - {file}")
    
    # Traiter les documents
    stats = process_documents()
    
    if not stats:
        print("\n❌ Erreur lors du traitement des documents.")
        exit(1)
    
    print("\n✅ Traitement terminé avec succès!")

